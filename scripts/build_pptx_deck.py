#!/usr/bin/env python3
"""Assemble the MetricLens figures into a clean 16:9 PPTX deck (EN + KR).

The four box-and-arrow diagrams (architecture, runtime, approach, model math) are
drawn with NATIVE PowerPoint objects -- rounded rectangles, text boxes, connector
arrows, freeform curves -- so every element is editable inside PowerPoint. Only the
brand logos are embedded as small images (rendered from Simple Icons). The five
evaluation charts are data plots and stay as PNG image slides.

Layout coordinates are ported verbatim from the SVG generators
(scripts/build_{architecture,runtime,approach,math}_diagram.py) via a px->slide
coordinate map, so the native slides match the published figures. NOTE: those SVG
scripts remain the source of truth; if a layout changes, update both.

Outputs (docs/): metriclens_deck.pptx (EN), metriclens_deck_kr.pptx (KR).
Usage:  python scripts/build_pptx_deck.py
"""

from __future__ import annotations

import math
import re
import urllib.request
from pathlib import Path

import cairosvg
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DIAG = ROOT / "docs" / "diagrams"
EVAL = ROOT / "docs" / "evaluation"
LOGO = DIAG / "_logos"
OUT = ROOT / "docs"
LOGO.mkdir(parents=True, exist_ok=True)

INK_H, SUB_H, ACC_H, RED_H = "#16191d", "#454b52", "#1f77b4", "#d62728"
INK, SUB, ACC = RGBColor(0x16, 0x19, 0x1D), RGBColor(0x45, 0x4B, 0x52), RGBColor(0x1F, 0x77, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SW, SH = Inches(13.333), Inches(7.5)  # 16:9 widescreen

DFONT = {"en": "Segoe UI", "kr": "Malgun Gothic"}   # normal text
MFONT = {"en": "Consolas", "kr": "Malgun Gothic"}    # equations / mono

# ----------------------------------------------------------------------------- logos
_PATH_RE = re.compile(r'<path[^>]*\sd="([^"]+)"')
_UA = "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 Chrome/124 Safari/537.36"
_SRC = ("https://cdn.jsdelivr.net/npm/simple-icons@13/icons/{slug}.svg",
        "https://cdn.simpleicons.org/{slug}")
SLUGS = ["react", "nginx", "fastapi", "python", "sqlalchemy", "postgresql",
         "sqlite", "googlecloud", "docker", "git"]


def fetch_path(slug):
    for tmpl in _SRC:
        try:
            req = urllib.request.Request(tmpl.format(slug=slug), headers={"User-Agent": _UA})
            with urllib.request.urlopen(req, timeout=20) as resp:
                m = _PATH_RE.search(resp.read().decode("utf-8"))
                if m:
                    return m.group(1)
        except Exception:  # noqa: BLE001
            continue
    print(f"  ! logo fetch failed for {slug}")
    return None


def logo_png(paths, slug, hexcolor):
    """Idempotently render a small single-colour PNG for ``slug``; None on failure."""
    d = paths.get(slug)
    if not d:
        return None
    fp = LOGO / f"{slug}_{hexcolor.lstrip('#')}.png"
    if not fp.exists():
        svg = (f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24">'
               f'<path d="{d}" fill="{hexcolor}"/></svg>')
        cairosvg.svg2png(bytestring=svg.encode("utf-8"), write_to=str(fp),
                         output_width=96, output_height=96)
    return fp


# ------------------------------------------------------------------- text-fit (px)
def _strwidth(s, size):
    return sum(size * (1.0 if ord(ch) >= 0x1100 else 0.52) for ch in s)


def _fit(s, maxw, size, minsize):
    while size > minsize and _strwidth(s, size) > maxw:
        size -= 0.5
    return size


def _hex(h):
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


# --------------------------------------------------------------- coordinate mapper
class Map:
    """Map an SVG canvas (Wsvg x Hsvg px) into a slide content box (inches)."""

    def __init__(self, wsvg, hsvg, cl, ct, cw, ch):
        self.s = min(cw / wsvg, ch / hsvg)            # inches per svg-px
        self.ox = cl + (cw - wsvg * self.s) / 2
        self.oy = ct + (ch - hsvg * self.s) / 2

    def X(self, px):
        return Inches(self.ox + px * self.s)

    def Y(self, px):
        return Inches(self.oy + px * self.s)

    def L(self, v):
        return Inches(v * self.s)

    def FS(self, px):
        return Pt(px * self.s * 72)


# ----------------------------------------------------------------- native helpers
def _no_autofit(tf):
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0


def rbox(slide, m, x, y, w, h, fill_h, line_h, lw_px=1.0, rx_px=8):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, m.X(x), m.Y(y), m.L(w), m.L(h))
    shp.adjustments[0] = max(0.0, min(0.5, rx_px / min(w, h)))
    shp.shadow.inherit = False
    if fill_h:
        shp.fill.solid(); shp.fill.fore_color.rgb = _hex(fill_h)
    else:
        shp.fill.background()
    if line_h:
        shp.line.color.rgb = _hex(line_h); shp.line.width = m.FS(lw_px)
    else:
        shp.line.fill.background()
    return shp


def _run(p, s, size_pt, color, font, bold=False, baseline=None):
    r = p.add_run(); r.text = s
    r.font.size = size_pt; r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    if baseline is not None:
        r._r.get_or_add_rPr().set("baseline", str(baseline))


def lines_box(slide, m, x, y, w, h, lines, font, align="l", anchor="m"):
    """lines: list of (text, size_px, hexcolor, bold) -> one paragraph each."""
    tb = slide.shapes.add_textbox(m.X(x), m.Y(y), m.L(w), m.L(h))
    tf = tb.text_frame; _no_autofit(tf)
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}[anchor]
    al = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al; p.space_before = Pt(0); p.space_after = Pt(0)
        _run(p, txt, m.FS(size), _hex(color), font, bold)
    return tb


def text1(slide, m, x, y, s, size_px, color_h, font, bold=False, align="l", width_px=520):
    """Single baseline-anchored SVG label -> slide textbox centred on y."""
    h = size_px * 1.7
    if align == "l":
        left = x
    elif align == "c":
        left = x - width_px / 2
    else:
        left = x - width_px
    tb = slide.shapes.add_textbox(m.X(left), Emu(int(m.Y(y - h / 2))), m.L(width_px), m.L(h))
    tf = tb.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    _run(p, s, m.FS(size_px), _hex(color_h), font, bold)
    return tb


def icon(slide, m, paths, slug, hexcolor, x, y, size):
    fp = logo_png(paths, slug, hexcolor)
    if fp:
        slide.shapes.add_picture(str(fp), m.X(x), m.Y(y), m.L(size), m.L(size))


def _style_line(cxn, m, color_h, lw_px, dash, head):
    ln = cxn.line; ln.color.rgb = _hex(color_h); ln.width = m.FS(lw_px)
    el = ln._get_or_add_ln()
    if dash:
        d = el.makeelement(qn("a:prstDash"), {"val": "dash"}); el.append(d)
    if head:
        t = el.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        el.append(t)


def connector(slide, m, x1, y1, x2, y2, color_h, lw_px=1.4, dash=False, head=True):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, m.X(x1), m.Y(y1), m.X(x2), m.Y(y2))
    _style_line(cxn, m, color_h, lw_px, dash, head)
    return cxn


def elbow(slide, m, pts, color_h, lw_px=1.4):
    """Orthogonal multi-segment arrow; head only on the final segment."""
    for i in range(len(pts) - 1):
        (x1, y1), (x2, y2) = pts[i], pts[i + 1]
        connector(slide, m, x1, y1, x2, y2, color_h, lw_px, head=(i == len(pts) - 2))


def curve(slide, m, pts, color_h, lw_px=1.4, dash=False, fill_h=None):
    fb = slide.shapes.build_freeform(int(m.X(pts[0][0])), int(m.Y(pts[0][1])), scale=1.0)
    fb.add_line_segments([(int(m.X(px)), int(m.Y(py))) for px, py in pts[1:]], close=bool(fill_h))
    shp = fb.convert_to_shape()
    shp.shadow.inherit = False
    if fill_h:
        shp.fill.solid(); shp.fill.fore_color.rgb = _hex(fill_h); shp.line.fill.background()
    else:
        shp.fill.background(); shp.line.color.rgb = _hex(color_h); shp.line.width = m.FS(lw_px)
        if dash:
            el = shp.line._get_or_add_ln()
            el.append(el.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return shp


_EQTOK = re.compile(r"_\{([^}]*)\}|\^\{([^}]*)\}|([^_^]+|[_^])")


def eq(slide, m, x, y, text, size_px, color_h, font, width_px=600):
    """Render an equation with _{sub}/^{sup} markup, baseline-anchored at y."""
    h = size_px * 1.7
    tb = slide.shapes.add_textbox(m.X(x), Emu(int(m.Y(y - h / 2))), m.L(width_px), m.L(h))
    tf = tb.text_frame; _no_autofit(tf); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    base, small = m.FS(size_px), m.FS(size_px * 0.72)
    col = _hex(color_h)
    for sub, sup, normal in _EQTOK.findall(text):
        if sub:
            _run(p, sub, small, col, font, baseline=-25000)
        elif sup:
            _run(p, sup, small, col, font, baseline=30000)
        else:
            _run(p, normal, base, col, font)
    return tb


# =============================================================== diagram: architecture
A_PRES = [("react", "#149ECA", "React 19 SPA", "Vite · ECharts (Canvas)", "Vite · ECharts (Canvas)"),
          ("nginx", "#009639", "nginx", "non-root static serving", "비루트 정적 서빙")]
A_BUS = [("fastapi", "#009688", "FastAPI", "Controller · Service · Repository", "Controller · Service · Repository"),
         ("python", "#3776AB", "Pure Core", "forecaster · optimizer", "forecaster · optimizer"),
         ("sqlalchemy", "#D71F00", "SQLAlchemy 2.0", "async ORM", "비동기 ORM")]
A_DATA = [("postgresql", "#4169E1", "PostgreSQL", "production (Cloud SQL)", "운영 (Cloud SQL)"),
          ("sqlite", "#003B57", "SQLite", "embedded demo (auto-seed)", "내장 데모 (자동 시드)")]
A_FLEET = [("googlecloud", "#4285F4", "Compute Engine", "VM inventory · setMachineType", "VM 인벤토리 · setMachineType"),
           ("googlecloud", "#4285F4", "Cloud Monitoring", "instance/cpu/utilization", "instance/cpu/utilization (CPU)"),
           ("googlecloud", "#4285F4", "Ops Agent", "memory/percent_used", "memory/percent_used (메모리)")]
A_RT = [("googlecloud", "#4285F4", "Cloud Run", "serverless · scale-to-zero", "서버리스 · scale-to-zero"),
        ("docker", "#2496ED", "Docker", "non-root images", "비루트 컨테이너 이미지")]
A_CI = [("git", "#F05032", "git push", "GitHub · main branch", "GitHub · main 브랜치"),
        ("googlecloud", "#4285F4", "Cloud Build", "lint·test·build → deploy", "린트·테스트·빌드 → 배포")]
A_TR = {
    "en": {"p": ("Presentation", "browser SPA on Cloud Run"),
           "b": ("Business", "layered FastAPI + dependency-free Core"),
           "d": ("Data", "managed Postgres (prod) · embedded SQLite (demo)"),
           "fleet": ("GCP managed fleet — external data plane", "live CPU/memory telemetry + resize actuation"),
           "ingest": "ingest · GcpSyncService: list_time_series → persist (Metric/Host)",
           "actuate": "actuate: setMachineType", "infra": "Deployment infrastructure (cross-cutting)",
           "rt": ("Runtime platform", "GCP-native, scale-to-zero"), "ci": ("CI/CD pipeline", "Cloud Build → Cloud Run"),
           "a1": "REST / HTTPS (CORS)", "a2": "SQLAlchemy 2.0 async"},
    "kr": {"p": ("표현 계층 (Presentation)", "브라우저 SPA · Cloud Run"),
           "b": ("업무 계층 (Business)", "레이어드 FastAPI + 의존성 없는 Core"),
           "d": ("데이터 계층 (Data)", "운영 Postgres · 데모 SQLite"),
           "fleet": ("GCP 관리 플릿 — 외부 데이터 플레인", "실시간 CPU/메모리 텔레메트리 + 리사이즈 실행"),
           "ingest": "수집 · GcpSyncService: list_time_series → 영속화(Metric/Host)",
           "actuate": "실행: setMachineType", "infra": "배포 인프라 (계층 횡단 관심사)",
           "rt": ("런타임 플랫폼", "GCP 네이티브, scale-to-zero"), "ci": ("CI/CD 파이프라인", "Cloud Build → Cloud Run"),
           "a1": "REST / HTTPS (CORS)", "a2": "SQLAlchemy 2.0 async"},
}


def _a_block(slide, m, paths, slug, color, name, role, x, y, w, h, font):
    rbox(slide, m, x, y, w, h, "#ffffff", "#16191d", 1.0, rx_px=8)
    icon(slide, m, paths, slug, color, x + 16, y + h / 2 - 14, 28)
    avail = w - 56 - 12
    lines_box(slide, m, x + 56, y, avail, h,
              [(name, _fit(name, avail, 14, 10.5), "#16191d", True),
               (role, _fit(role, avail, 11, 8.5), "#454b52", False)], font, "l", "m")


def _a_band(slide, m, paths, label, sub, items, y, W, lang, font, h=88, bw=280):
    left_col = 230; area = (W - 56) - left_col; gap = 26; bh = 64
    n = len(items); total = n * bw + (n - 1) * gap; startx = left_col + (area - total) / 2
    lw = startx - 60 - 14
    rbox(slide, m, 40, y, W - 80, h, "#ffffff", "#16191d", 1.0, rx_px=10)
    lines_box(slide, m, 60, y, lw, h,
              [(label, _fit(label, lw, 14, 9), "#16191d", True),
               (sub, _fit(sub, lw, 10, 8), "#454b52", False)], font, "l", "m")
    by = y + h / 2 - bh / 2
    for i, (slug, color, name, role_en, role_kr) in enumerate(items):
        role = role_en if lang == "en" else role_kr
        _a_block(slide, m, paths, slug, color, name, role, startx + i * (bw + gap), by, bw, bh, font)


def _a_panel(slide, m, paths, label, sub, items, x, y, w, h, lang, font):
    rbox(slide, m, x, y, w, h, "#ffffff", "#16191d", 1.0, rx_px=10)
    lines_box(slide, m, x + 18, y + 8, w - 36, 40,
              [(label, 13, "#16191d", True), (sub, 10, "#454b52", False)], font, "l", "t")
    pad = 18; gap = 18; bh = 64; n = len(items); bw = (w - 2 * pad - (n - 1) * gap) / n
    by = y + h - pad - bh
    for i, (slug, color, name, role_en, role_kr) in enumerate(items):
        role = role_en if lang == "en" else role_kr
        _a_block(slide, m, paths, slug, color, name, role, x + pad + i * (bw + gap), by, bw, bh, font)


def draw_architecture(slide, m, lang, paths):
    T = A_TR[lang]; W = 1280; cx = W / 2; f = DFONT[lang]
    _a_band(slide, m, paths, *T["p"], A_PRES, 44, W, lang, f)
    connector(slide, m, cx, 132, cx, 162, "#16191d")
    text1(slide, m, cx + 14, 147, T["a1"], 11, SUB_H, f, width_px=240)
    _a_band(slide, m, paths, *T["b"], A_BUS, 162, W, lang, f)
    connector(slide, m, cx, 250, cx, 280, "#16191d")
    text1(slide, m, cx + 14, 265, T["a2"], 11, SUB_H, f, width_px=240)
    _a_band(slide, m, paths, *T["d"], A_DATA, 280, W, lang, f)
    _a_band(slide, m, paths, *T["fleet"], A_FLEET, 424, W, lang, f)
    # control loop in the 56px channel
    connector(slide, m, 380, 424, 380, 368, ACC_H)            # ingest (up)
    text1(slide, m, 394, 400, T["ingest"], 11, SUB_H, f, width_px=520)
    connector(slide, m, 900, 368, 900, 424, RED_H)            # actuate (down)
    text1(slide, m, 914, 396, T["actuate"], 11, SUB_H, f, width_px=300)
    # cross-cutting infra
    text1(slide, m, 40, 536, T["infra"], 12.5, INK_H, f, bold=True, width_px=400)
    connector(slide, m, 40, 550, W - 40, 550, "#16191d", 1.0, dash=True, head=False)
    _a_panel(slide, m, paths, *T["rt"], A_RT, 40, 560, 590, 128, lang, f)
    _a_panel(slide, m, paths, *T["ci"], A_CI, 650, 560, 590, 128, lang, f)


# =============================================================== diagram: runtime
R_TR = {
    "en": {"b1t": "CI/CD pipeline (Cloud Build to Cloud Run)",
           "b1s": "git push -> lint -> test -> build image -> deploy -> health check",
           "dev": "Developer", "build_sub": "lint·test·build·deploy", "ar_sub": "container images",
           "b2t": "Runtime serving (serverless)", "b2s": "browser -> frontend -> backend; secrets from Secret Manager",
           "browser": "Browser", "frontend": "Frontend", "backend": "Backend", "sm_sub": "DB DSN / secrets",
           "b3t": "Real-fleet monitoring & resize (Cloud Monitoring + Compute Engine)",
           "b3s": "backend reads live CPU/memory and resizes real VMs within the budget guard",
           "mon_sub": "CPU + memory", "fleet": "Real fleet · ml-web/api/batch/idle",
           "fleet_sub": "e2-small · load generator + Ops Agent", "trigger": "trigger", "image": "image",
           "deploy": "deploy", "secrets": "secrets", "read": "read CPU/mem", "resize": "resize (budget-guarded)",
           "scs": "stop->change->start", "push": "push metrics (Ops Agent)"},
    "kr": {"b1t": "CI/CD 파이프라인 (Cloud Build → Cloud Run)",
           "b1s": "git push → 린트 → 테스트 → 이미지 빌드 → 배포 → 헬스 체크",
           "dev": "개발자(Developer)", "build_sub": "린트·테스트·빌드·배포", "ar_sub": "컨테이너 이미지",
           "b2t": "런타임 서빙 (서버리스)", "b2s": "브라우저 → 프론트엔드 → 백엔드, 비밀값은 Secret Manager",
           "browser": "브라우저(Browser)", "frontend": "프론트엔드", "backend": "백엔드", "sm_sub": "DB DSN / 비밀값",
           "b3t": "실제 플릿 모니터링 · 리사이즈 (Cloud Monitoring + Compute Engine)",
           "b3s": "백엔드가 실측 CPU/메모리를 읽고, 예산 한도 내에서 실제 VM을 리사이즈",
           "mon_sub": "CPU + 메모리", "fleet": "실제 플릿 · ml-web/api/batch/idle",
           "fleet_sub": "e2-small · 부하 생성기 + Ops Agent", "trigger": "트리거", "image": "이미지",
           "deploy": "배포", "secrets": "비밀값", "read": "CPU/메모리 조회", "resize": "리사이즈 (예산 가드)",
           "scs": "stop→change→start", "push": "메트릭 전송 (Ops Agent)"},
}


def _r_band(slide, m, x, y, w, h, label, sub, font):
    rbox(slide, m, x, y, w, h, "#ffffff", "#16191d", 1.0, rx_px=10)
    lines_box(slide, m, x + 16, y + 8, w - 32, 38,
              [(label, 13, "#16191d", True), (sub, 10, "#454b52", False)], font, "l", "t")


def _r_box(slide, m, paths, x, y, w, h, slug, color, title, sub, font):
    rbox(slide, m, x, y, w, h, "#ffffff", "#16191d", 1.0, rx_px=8)
    tx = x + 14
    if slug:
        icon(slide, m, paths, slug, color, tx, y + h / 2 - 13, 26); tx += 36
    avail = (x + w) - tx - 10
    lines_box(slide, m, tx, y, avail, h,
              [(title, _fit(title, avail, 13, 9), "#16191d", True),
               (sub, _fit(sub, avail, 10.5, 8), "#454b52", False)], font, "l", "m")


def draw_runtime(slide, m, lang, paths):
    T = R_TR[lang]; W = 1140; f = DFONT[lang]; bh = 58; G = "#4285F4"
    _r_band(slide, m, 30, 40, W - 60, 110, T["b1t"], T["b1s"], f)
    y1 = 86
    _r_box(slide, m, paths, 50, y1, 150, bh, "git", "#F05032", T["dev"], "git push", f)
    _r_box(slide, m, paths, 270, y1, 230, bh, "googlecloud", G, "Cloud Build", T["build_sub"], f)
    _r_box(slide, m, paths, 560, y1, 180, bh, "googlecloud", G, "Artifact Registry", T["ar_sub"], f)
    _r_box(slide, m, paths, 800, y1, 190, bh, "googlecloud", G, "Cloud Run", "scale-to-zero", f)
    for a, b, lab in ((200, 270, T["trigger"]), (500, 560, T["image"]), (740, 800, T["deploy"])):
        connector(slide, m, a, y1 + bh / 2, b, y1 + bh / 2, "#16191d")
        text1(slide, m, (a + b) / 2, y1 + bh / 2 - 12, lab, 10.5, SUB_H, f, align="c", width_px=120)

    _r_band(slide, m, 30, 190, W - 60, 110, T["b2t"], T["b2s"], f)
    y2 = 236
    _r_box(slide, m, paths, 50, y2, 150, bh, None, "", T["browser"], "React SPA", f)
    _r_box(slide, m, paths, 270, y2, 190, bh, "googlecloud", G, T["frontend"], "Cloud Run · nginx", f)
    _r_box(slide, m, paths, 520, y2, 190, bh, "googlecloud", G, T["backend"], "FastAPI · Cloud Run", f)
    _r_box(slide, m, paths, 800, y2, 190, bh, "googlecloud", G, "Secret Manager", T["sm_sub"], f)
    for a, b, lab in ((200, 270, "HTTPS"), (460, 520, "REST"), (710, 800, T["secrets"])):
        connector(slide, m, a, y2 + bh / 2, b, y2 + bh / 2, "#16191d")
        text1(slide, m, (a + b) / 2, y2 + bh / 2 - 12, lab, 10.5, SUB_H, f, align="c", width_px=120)

    _r_band(slide, m, 30, 340, W - 60, 240, T["b3t"], T["b3s"], f)
    _r_box(slide, m, paths, 80, 418, 210, 56, "googlecloud", G, "Cloud Monitoring", T["mon_sub"], f)
    _r_box(slide, m, paths, 470, 418, 230, 56, "googlecloud", G, "Compute Engine API", "setMachineType", f)
    _r_box(slide, m, paths, 820, 418, 270, 68, None, "", T["fleet"], T["fleet_sub"], f)
    be_b = 294; mon_cx = 185
    elbow(slide, m, [(615, be_b), (615, 414)], "#16191d")
    text1(slide, m, 625, 352, T["resize"], 10.5, SUB_H, f, width_px=200)
    elbow(slide, m, [(560, be_b), (560, 402), (mon_cx, 402), (mon_cx, 414)], "#16191d")
    text1(slide, m, 378, 396, T["read"], 10.5, SUB_H, f, align="c", width_px=160)
    elbow(slide, m, [(700, 446), (816, 446)], "#16191d")
    text1(slide, m, 760, 438, T["scs"], 10.5, SUB_H, f, align="c", width_px=160)
    elbow(slide, m, [(955, 486), (955, 540), (mon_cx, 540), (mon_cx, 478)], "#16191d")
    text1(slide, m, 575, 534, T["push"], 10.5, SUB_H, f, align="c", width_px=220)


# =============================================================== diagram: approach
AP_TR = {
    "en": {"cards": [("Metric history", ["CPU · memory · network", "time series (hourly)"]),
                     ("Decomposition + AR", ["trend + seasonal", "+ rho·residual -> forecast"]),
                     ("Robust peak (p95)", ["p95 of the forecast", "(ignores transient spikes)"]),
                     ("Integer program", ["smallest allocation", "s.t. headroom (exact)"]),
                     ("GCP machine type", ["snap to the nearest", "E2/N2/C2/C3 instance"]),
                     ("Resize + audit", ["SLO preserved", "persisted to audit log"])],
           "sketch": "forecast + 95% interval", "eq": "peak_load × safety_margin  ≤  target_utilisation × allocation",
           "f2": "exact enumeration of the smallest feasible integer size"},
    "kr": {"cards": [("메트릭 이력", ["CPU · 메모리 · 네트워크", "시계열 (시간단위)"]),
                     ("분해 + AR 보정", ["추세(trend) + 계절(seasonal)", "+ rho·잔차 → 예측(forecast)"]),
                     ("강건 피크 (p95)", ["예측의 p95 통계", "(단발 스파이크 무시)"]),
                     ("정수계획 최적화", ["최소 할당 탐색", "s.t. 헤드룸 (정확 해)"]),
                     ("GCP 머신 타입", ["가장 근접한", "E2/N2/C2/C3 인스턴스"]),
                     ("리사이즈 + 감사", ["SLO 보존", "감사 로그 영속 기록"])],
           "sketch": "예측(forecast) + 95% 구간(interval)", "eq": "peak_load × safety_margin  ≤  target_utilisation × allocation",
           "f2": "전수 열거(exact enumeration)로 최소 할당 탐색 — 실현 가능한 최소 정수 사양"},
}


def _ap_card(slide, m, n, x, y, w, h, title, lines, font):
    rbox(slide, m, x, y, w, h, "#ffffff", "#16191d", 1.0, rx_px=9)
    block_h = 26 + len(lines) * 18; top = y + (h - block_h) / 2
    # number disc
    disc = slide.shapes.add_shape(MSO_SHAPE.OVAL, m.X(x + 22 - 13), m.Y(top + 8 - 13), m.L(26), m.L(26))
    disc.shadow.inherit = False; disc.fill.solid(); disc.fill.fore_color.rgb = ACC; disc.line.fill.background()
    dtf = disc.text_frame; _no_autofit(dtf); dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    dp = dtf.paragraphs[0]; dp.alignment = PP_ALIGN.CENTER
    _run(dp, str(n), m.FS(13), WHITE, font, bold=True)
    text1(slide, m, x + 44, top + 13, title, _fit(title, w - 58, 14, 9.5), INK_H, font, bold=True, width_px=w - 50)
    for i, ln in enumerate(lines):
        text1(slide, m, x + 16, top + 40 + i * 18, ln, _fit(ln, w - 28, 11.5, 8.5), SUB_H, font, width_px=w - 24)


def _ap_sketch(slide, m, x, y, w, h, label, font):
    rbox(slide, m, x - 8, y - 12, w + 16, h + 24, "#ffffff", "#16191d", 1.0, rx_px=6)
    text1(slide, m, x - 2, y - 18, label, 8.5, SUB_H, font, width_px=240)
    n, split = 28, 18
    pts = [(x + (i / (n - 1)) * w,
            y + h - (0.5 + 0.32 * math.sin((i / (n - 1)) * 6.0) - 0.12 * (i / (n - 1))) * h)
           for i in range(n)]
    band = ([(px, py - 9) for px, py in pts[split - 1:]]
            + [(pts[i][0], pts[i][1] + 9) for i in range(n - 1, split - 2, -1)])
    curve(slide, m, band, None, fill_h="#d6e6f3")
    curve(slide, m, pts[:split], "#111111", 1.4)
    curve(slide, m, pts[split - 1:], "#1f77b4", 1.4, dash=True)


def draw_approach(slide, m, lang):
    T = AP_TR[lang]; W = 1280; f = DFONT[lang]; mf = MFONT[lang]
    cards = T["cards"]; n = len(cards); cw, ch, gap = 184, 132, 20
    x0 = (W - (n * cw + (n - 1) * gap)) / 2; cy = 350
    for i, (title, lines) in enumerate(cards):
        x = x0 + i * (cw + gap)
        _ap_card(slide, m, i + 1, x, cy, cw, ch, title, lines, f)
        if i:
            connector(slide, m, x - gap, cy + ch / 2, x, cy + ch / 2, "#16191d")
    c2x = x0 + (cw + gap); sk_w = cw + 36
    _ap_sketch(slide, m, c2x + (cw - sk_w) / 2 + 8, 120, sk_w - 16, 150, T["sketch"], f)
    connector(slide, m, c2x + cw / 2, 282, c2x + cw / 2, cy, "#16191d", 1.0, dash=True, head=False)
    fw, fh = 820, 104; fx = (W - fw) / 2; fy = cy + ch + 58
    rbox(slide, m, fx, fy, fw, fh, "#ffffff", "#1f77b4", 1.0, rx_px=10)
    eqs = _fit(T["eq"], fw - 48, 16, 10)
    text1(slide, m, fx + fw / 2, fy + 44, T["eq"], eqs, INK_H, mf, bold=True, align="c", width_px=fw - 40)
    text1(slide, m, fx + fw / 2, fy + 74, T["f2"], _fit(T["f2"], fw - 60, 12, 9), SUB_H, f, align="c", width_px=fw - 50)
    connector(slide, m, fx + fw / 2, fy, x0 + 3 * (cw + gap) + cw / 2, cy + ch, "#1f77b4", 1.0, dash=True, head=False)


# =============================================================== diagram: model math
M_ACC_F, M_ACC_O = "#1f77b4", "#d62728"
M_BAND = "#f4f6f8"
M_TR = {
    "en": {"model": "Decomposition model:",
           "laneF": "A. Forecasting  —  seasonal-trend decomposition + AR(1) residual correction  (CPU-only, stdlib)",
           "laneO": "B. Resizing  —  exact integer program over a bounded space (per resource: vCPU, memory blocks)",
           "fc": [("Trend (OLS)", ["T_{t} = b.t + a", "block means y_{k},", "centres c_{k} (deseason)"]),
                  ("Seasonal index", ["d_{t} = y_{t} - T_{t}", "S_{j}=mean d (t mod m=j)", "re-centred: sum S_{j}=0"]),
                  ("Residual + AR(1)", ["r_{t} = d_{t} - S_{t mod m}", "rho = Sr_{t}r_{t-1} / Sr_{t}^{2}", "rho in [0, 0.95]"]),
                  ("Point forecast", ["y^_{n-1+h} = T_{n-1+h}", " + S_{(n-1+h) mod m}", " + rho^{h}.r_{n-1}"]),
                  ("95% interval", ["RMSE: 1-step backtest", "PI = y^ +/- 1.96.RMSE", "PICP ~ 0.93-0.98"])],
           "op": [("Robust peak", ["x = p95 of forecast", "rank = ceil(0.95.N)", "(nearest-rank, not max)"]),
                  ("Load (in units)", ["L = (x/100).u_{cur}.g", "g = 1.2 (safety margin)", "u_{cur}: current units"]),
                  ("Integer program", ["minimise  u", "s.t.  L <= t.u", "t=0.65, u in {1..u_cur}"]),
                  ("Allocation + cost", ["u* = min(u_{cur}, ceil(L/t))", "M* = ceil(L_{m}/t).256MB", "save = avg D/cur .100%"])],
           "feed": "forecast peak feeds the optimizer",
           "out": "-> snap to GCP machine type (E2 / N2 / C2 / C3), budget-guarded resize", "mfont": "Consolas"},
    "kr": {"model": "분해 모델:",
           "laneF": "A. 예측  —  계절-추세 분해 + AR(1) 잔차 보정  (CPU 전용, 표준 라이브러리)",
           "laneO": "B. 리사이징  —  유계 공간 정확 정수계획 (자원별: vCPU, 메모리 블록)",
           "fc": [("추세 (OLS)", ["T_{t} = b.t + a", "블록평균 y_{k},", "중심 c_{k} (계절제거)"]),
                  ("계절 지수", ["d_{t} = y_{t} - T_{t}", "S_{j}=mean d (t mod m=j)", "재중심화: sum S_{j}=0"]),
                  ("잔차 + AR(1)", ["r_{t} = d_{t} - S_{t mod m}", "rho = Sr_{t}r_{t-1} / Sr_{t}^{2}", "rho in [0, 0.95]"]),
                  ("점예측", ["y^_{n-1+h} = T_{n-1+h}", " + S_{(n-1+h) mod m}", " + rho^{h}.r_{n-1}"]),
                  ("95% 예측구간", ["RMSE: 1-스텝 백테스트", "PI = y^ +/- 1.96.RMSE", "PICP ~ 0.93-0.98"])],
           "op": [("강건 피크", ["x = 예측의 p95", "rank = ceil(0.95.N)", "(최댓값 아닌 근접순위)"]),
                  ("부하(자원단위)", ["L = (x/100).u_{cur}.g", "g = 1.2 (안전마진)", "u_{cur}: 현재 단위수"]),
                  ("정수계획", ["minimise  u", "s.t.  L <= t.u", "t=0.65, u in {1..u_cur}"]),
                  ("할당 + 절감", ["u* = min(u_{cur}, ceil(L/t))", "M* = ceil(L_{m}/t).256MB", "절감 = avg D/cur .100%"])],
           "feed": "예측 피크가 옵티마이저 입력",
           "out": "-> GCP 머신 타입(E2 / N2 / C2 / C3)으로 스냅, 예산 가드 리사이즈", "mfont": "Malgun Gothic"},
}


def _m_card(slide, m, x, y, w, h, n, title, eqs, accent, font, mf):
    rbox(slide, m, x, y, w, h, "#ffffff", "#16191d", 1.0, rx_px=9)
    rbox(slide, m, x, y, w, 26, accent, None, rx_px=9)            # rounded header
    rbox(slide, m, x, y + 14, w, 12, accent, None, rx_px=0)        # square off the bottom of header
    text1(slide, m, x + 11, y + 13, f"{n}. {title}", 11.5, "#ffffff", font, bold=True, width_px=w - 16)
    ey = y + 48
    for line in eqs:
        eq(slide, m, x + 12, ey, line, 12.5, "#16191d", mf, width_px=w - 18)
        ey += 22


def draw_math(slide, m, lang):
    T = M_TR[lang]; W = 1180; cx = W / 2; f = DFONT[lang]; mf = MFONT[lang]
    rbox(slide, m, 40, 44, W - 80, 34, M_BAND, "#16191d", 1.0, rx_px=8)
    text1(slide, m, 56, 61, T["model"], 12, SUB_H, f, bold=True, width_px=240)
    eq(slide, m, 190, 66, "y_{t} = T_{t} + S_{t mod m} + r_{t}", 14.5, INK_H, mf, width_px=420)
    text1(slide, m, W - 56, 61, "m = period", 11, SUB_H, f, align="r", width_px=160)

    text1(slide, m, 40, 104, T["laneF"], 12.5, M_ACC_F, f, bold=True, width_px=1100)
    n, cw, gap, cy, ch = 5, 200, 24, 116, 132
    x0 = (W - (n * cw + (n - 1) * gap)) / 2
    for i, (title, eqs) in enumerate(T["fc"]):
        x = x0 + i * (cw + gap)
        _m_card(slide, m, x, cy, cw, ch, i + 1, title, eqs, M_ACC_F, f, mf)
        if i:
            connector(slide, m, x - gap, cy + ch / 2, x, cy + ch / 2, M_ACC_F)
    feed_y = cy + ch
    connector(slide, m, cx, feed_y, cx, feed_y + 40, M_ACC_O)
    text1(slide, m, cx + 12, feed_y + 22, T["feed"], 11, SUB_H, f, width_px=300)

    ly = feed_y + 64
    text1(slide, m, 40, ly, T["laneO"], 12.5, M_ACC_O, f, bold=True, width_px=1100)
    n2, cw2, gap2, cy2 = 4, 254, 26, ly + 12
    x0b = (W - (n2 * cw2 + (n2 - 1) * gap2)) / 2
    for i, (title, eqs) in enumerate(T["op"]):
        x = x0b + i * (cw2 + gap2)
        _m_card(slide, m, x, cy2, cw2, ch, i + 6, title, eqs, M_ACC_O, f, mf)
        if i:
            connector(slide, m, x - gap2, cy2 + ch / 2, x, cy2 + ch / 2, M_ACC_O)
    text1(slide, m, cx, cy2 + ch + 30, T["out"], 12.5, INK_H, f, bold=True, align="c", width_px=900)


# ----------------------------------------------------------------- deck assembly
TITLE = {
    "en": ("MetricLens AI",
           "Lightweight time-series forecasting + integer-programming resizing for SLO-safe minimal allocation",
           "Team Googling — Lee Yongwon · Seo Hannyeong · Song Simwoo"),
    "kr": ("MetricLens AI", "경량 시계열 예측 + 정수계획 리사이징으로 SLO 보장 최소 자원 산출",
           "Team 구글링 — 이용원 · 서한녕 · 송심우"),
}
# (title, subtitle, kind, payload). kind: "native"->(fn, Wsvg, Hsvg) ; "image"->basename
SLIDES = {
    "en": [
        ("System architecture", "3 logical layers + live GCP data plane · Cloud Run", "native", (draw_architecture, 1280, 720)),
        ("Runtime & CI/CD", "Cloud Build -> Cloud Run · Monitoring ingest / Compute Engine resize loop", "native", (draw_runtime, 1140, 620)),
        ("Method pipeline", "telemetry -> forecast -> SLO-constrained integer program -> resize", "native", (draw_approach, 1280, 720)),
        ("Model mathematics", "STL + AR(1) forecast and integer-programming resize: the exact equations", "native", (draw_math, 1180, 660)),
        ("Forecast accuracy: predicted vs actual", "one-step forecast with the 95% prediction interval", "image", "fig_forecast_overlay"),
        ("Forecast error (RMSE)", "model vs last-value / seasonal-naive baselines (lower is better)", "image", "fig_error_rmse"),
        ("MASE", "below 1 beats the seasonal-naive baseline", "image", "fig_mase"),
        ("Prediction-interval calibration (PICP)", "empirical coverage 0.93-0.98 vs nominal 0.95", "image", "fig_coverage"),
        ("Residual autocorrelation (ACF)", "within the white-noise band = good fit", "image", "fig_residual_acf"),
    ],
    "kr": [
        ("시스템 아키텍처", "3개 논리 계층 + 실시간 GCP 데이터 플레인 · Cloud Run", "native", (draw_architecture, 1280, 720)),
        ("런타임 · CI/CD 구동", "Cloud Build → Cloud Run · 실측 적재/Compute Engine 리사이즈 루프", "native", (draw_runtime, 1140, 620)),
        ("방법론 파이프라인", "텔레메트리 → 예측 → SLO 제약 정수계획 → 리사이즈", "native", (draw_approach, 1280, 720)),
        ("모델 수식 연산", "STL + AR(1) 예측과 정수계획 리사이징의 실제 계산식", "native", (draw_math, 1180, 660)),
        ("예측 정확도: 예측 vs 실측", "1-스텝 예측과 95% 예측구간 (대표 워크로드)", "image", "fig_forecast_overlay"),
        ("예측 오차 (RMSE)", "모델 vs 직전값/seasonal-naive 기준선 (낮을수록 우수)", "image", "fig_error_rmse"),
        ("MASE", "1 미만이면 seasonal-naive 기준선 능가", "image", "fig_mase"),
        ("예측구간 보정 (PICP)", "실측 커버리지 0.93–0.98 vs 공칭 0.95", "image", "fig_coverage"),
        ("잔차 자기상관 (ACF)", "백색잡음 대역 내 = 적합 양호", "image", "fig_residual_acf"),
    ],
}


def _slide_header(slide, title, subtitle, font):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = tb.text_frame; _no_autofit(tf)
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.name = font; r.font.color.rgb = INK
    tb2 = slide.shapes.add_textbox(Inches(0.52), Inches(1.02), Inches(12.3), Inches(0.45))
    tf2 = tb2.text_frame; _no_autofit(tf2)
    r2 = tf2.paragraphs[0].add_run(); r2.text = subtitle
    r2.font.size = Pt(14); r2.font.name = font; r2.font.color.rgb = SUB
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.045))
    bar.shadow.inherit = False; bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background()


def build_one(lang, suf, paths):
    f = DFONT[lang]
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # title slide
    s = prs.slides.add_slide(blank)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SW, Inches(0.06))
    bar.shadow.inherit = False; bar.fill.solid(); bar.fill.fore_color.rgb = ACC; bar.line.fill.background()
    name, sub, team = TITLE[lang]
    for txt, top, size, color, bold in ((name, 2.7, 48, INK, True), (sub, 3.95, 20, SUB, False),
                                        (team, 6.4, 15, ACC, False)):
        tb = s.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.5), Inches(1.2))
        tf = tb.text_frame; tf.word_wrap = True
        r = tf.paragraphs[0].add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = f; r.font.color.rgb = color

    total = len(SLIDES[lang]) + 1
    for idx, (title, subtitle, kind, payload) in enumerate(SLIDES[lang], start=2):
        s = prs.slides.add_slide(blank)
        _slide_header(s, title, subtitle, f)
        if kind == "native":
            fn, wsvg, hsvg = payload
            m = Map(wsvg, hsvg, 0.3, 1.5, 12.73, 5.7)
            fn(s, m, lang, paths) if fn in (draw_architecture, draw_runtime) else fn(s, m, lang)
        else:
            img = EVAL / f"{payload}{suf}.png"
            if img.exists():
                iw, ih = Image.open(img).size; ar = iw / ih
                bl, bt, bw, bh = 0.5, 1.7, 12.33, 5.35
                if ar >= bw / bh:
                    w = bw; h = bw / ar
                else:
                    h = bh; w = bh * ar
                s.shapes.add_picture(str(img), Inches(bl + (bw - w) / 2), Inches(bt + (bh - h) / 2),
                                     Inches(w), Inches(h))
        pn = s.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.4))
        rp = pn.text_frame.paragraphs[0]; rp.alignment = PP_ALIGN.RIGHT
        r = rp.add_run(); r.text = f"{idx}/{total}"; r.font.size = Pt(10); r.font.name = f; r.font.color.rgb = SUB

    dest = OUT / f"metriclens_deck{suf}.pptx"
    prs.save(str(dest))
    print(f"Wrote {dest} ({total} slides)")


def build():
    print("Fetching official OSS logos…")
    paths = {s: fetch_path(s) for s in SLUGS}
    build_one("en", "", paths)
    build_one("kr", "_kr", paths)


if __name__ == "__main__":
    build()
